from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "Trustnode_edge_app" / "docs"
ASSETS = DOCS / "assets"
LOGO = ROOT / "trustnode_logo.png"
OUT = DOCS / "TRUSTNODE_SCALABILITY_PITCH_DECK_2026-04-22.pptx"

NAVY = RGBColor(13, 19, 35)
NAVY2 = RGBColor(21, 40, 72)
BLUE = RGBColor(0, 133, 255)
CYAN = RGBColor(34, 211, 238)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(246, 249, 253)
TXT = RGBColor(18, 30, 49)
MUTED = RGBColor(81, 99, 125)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def topbar(slide, title, dark=True):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY if dark else NAVY2
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "  " + title
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE


def footer(slide, text="TrustNode Pitch Deck"):
    f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.28), prs.slide_width, Inches(0.28))
    f.fill.solid()
    f.fill.fore_color.rgb = RGBColor(232, 239, 248)
    f.line.fill.background()
    tf = f.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text + "  "
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def logo(slide, x=11.7, y=0.1, h=0.4):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), height=Inches(h))


def title_block(slide, title, subtitle):
    t = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(8.5), Inches(2.2))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(42)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(190, 231, 255)


def bullets(slide, x, y, w, h, items, color=TXT, fs=20):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.space_after = Pt(10)


def card(slide, x, y, w, h, title, body, accent=BLUE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = RGBColor(211, 223, 238)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    t = slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.2), Inches(w-0.36), Inches(h-0.3))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = NAVY2
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(13)
    p2.font.color.rgb = MUTED


def image_box(slide, img, x, y, w, h):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x-0.05), Inches(y-0.05), Inches(w+0.10), Inches(h+0.10))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(206, 220, 239)
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


# 1 Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
left = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.8), Inches(7.5))
left.fill.solid(); left.fill.fore_color.rgb = NAVY2; left.line.fill.background()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.7), Inches(13.333), Inches(0.8))
accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background()
logo(s, x=0.75, y=0.55, h=0.8)
title_block(s, "TrustNode", "Scalable Industrial Edge-to-Cloud Platform")
bullets(s, 5.3, 3.3, 7.3, 2.8, [
    "Real-time PLC + meter telemetry",
    "Secure multi-customer cloud architecture",
    "Low-cost, fast rollout, enterprise-ready path",
], color=WHITE, fs=18)

# 2 Opportunity
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "Why This Matters"); logo(s); footer(s)
card(s, 0.7, 1.1, 3.9, 2.3, "Plant challenge", "Data is fragmented across OT systems and hard to access remotely.", ORANGE)
card(s, 4.8, 1.1, 3.9, 2.3, "Business impact", "Limited visibility delays decisions and increases downtime/energy waste.", BLUE)
card(s, 8.9, 1.1, 3.7, 2.3, "Need", "Secure real-time access with audit-grade data integrity.", CYAN)
card(s, 0.7, 3.8, 12.0, 2.7, "TrustNode value", "Single architecture for local edge operations and cloud client access, preserving exact telemetry provenance and enabling commercial multi-tenant delivery.", GREEN)

# 3 Solution at glance
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "Solution At A Glance"); logo(s); footer(s)
bullets(s, 0.9, 1.2, 6.0, 4.8, [
    "Edge-first durability",
    "Outbox sync with retries",
    "Tenant-isolated cloud ingest",
    "Live dashboards from latest state",
    "Historian/reporting from immutable raw data",
], color=TXT, fs=20)
card(s, 6.8, 1.2, 5.8, 5.3, "Commercial-ready with minimal change", "Keep current architecture, add control plane for customer onboarding, licensing, edge activation, and URL tenancy.", BLUE)

# 4 Target architecture diagram
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "Target Architecture"); logo(s); footer(s)
image_box(s, ASSETS / "topology_2.png", 0.7, 1.0, 12.0, 5.9)

# 5 Security architecture
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "Security and Data Isolation"); logo(s); footer(s)
image_box(s, ASSETS / "topology_3.png", 0.7, 1.0, 12.0, 3.6)
card(s, 0.7, 4.85, 3.9, 1.7, "OT safety", "Outbound-only edge traffic. No inbound internet to OT.", GREEN)
card(s, 4.85, 4.85, 3.9, 1.7, "Tenant protection", "Token scope + API checks + RLS for user-facing reads.", BLUE)
card(s, 9.0, 4.85, 3.7, 1.7, "Admin security", "MFA + security audit log for privileged actions.", ORANGE)

# 6 Scale model
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "Scalability Model"); logo(s); footer(s)
card(s, 0.7, 1.2, 3.9, 2.5, "Tier A", "Shared DB + RLS\nLowest cost\nDefault for SMB", BLUE)
card(s, 4.8, 1.2, 3.9, 2.5, "Tier B", "Shared cluster + isolated schema\nHigher isolation", CYAN)
card(s, 8.9, 1.2, 3.8, 2.5, "Tier C", "Dedicated project/cluster\nEnterprise compliance", GREEN)
card(s, 0.7, 4.1, 12.0, 2.3, "Performance targets", "Local UI <= 1s, Cloud UI <= 2s p95, deterministic outbox sync, full provenance for historian/reporting.", BLUE)

# 7 Rollout
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "90-Day Rollout Plan"); logo(s); footer(s)
for i, (title, txt) in enumerate([
    ("30 days", "Stabilize ingest bootstrap + diagnostics"),
    ("60 days", "Control plane core + activation + licensing"),
    ("90 days", "Wildcard customer URLs + MFA + isolation tests"),
]):
    x = 0.9 + i * 4.2
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3.7), Inches(3.9))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.color.rgb = RGBColor(205, 219, 237)
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+1.35), Inches(1.1), Inches(1.0), Inches(1.0))
    n.fill.solid(); n.fill.fore_color.rgb = BLUE; n.line.fill.background()
    nt = n.text_frame; nt.clear(); np = nt.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
    nr = np.add_run(); nr.text = str(i+1); nr.font.size = Pt(26); nr.font.bold = True; nr.font.color.rgb = WHITE
    t = s.shapes.add_textbox(Inches(x+0.2), Inches(2.35), Inches(3.3), Inches(2.6)).text_frame
    t.clear()
    p = t.paragraphs[0]; p.text = title; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY2
    p2 = t.add_paragraph(); p2.text = txt; p2.font.size = Pt(14); p2.font.color.rgb = MUTED

# 8 Commercial model
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG); topbar(s, "Commercial Packaging"); logo(s); footer(s)
card(s, 0.7, 1.2, 5.9, 2.1, "Core package", "Edge collector + cloud client + dashboards + historian + reporting", BLUE)
card(s, 6.9, 1.2, 5.8, 2.1, "Add-on modules", "Power management, advanced analytics, AI advisory insights", CYAN)
card(s, 0.7, 3.6, 5.9, 2.8, "Recurring value", "Remote operations visibility, lower downtime, faster diagnostics, secure multi-site rollout", GREEN)
card(s, 6.9, 3.6, 5.8, 2.8, "Upsell path", "Isolation tiers, dedicated tenancy, compliance bundle, AI optimization", ORANGE)

# 9 Closing
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY); topbar(s, "Next Step"); logo(s); footer(s, "TrustNode - Investor/Demo Deck")
t = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(1.0)).text_frame
p = t.paragraphs[0]; r = p.add_run(); r.text = "Ready for Pilot-to-Scale Execution"; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
bullets(s, 0.95, 2.8, 11.6, 3.2, [
    "Deploy control-plane Phase 0 now",
    "Run 2-3 customer pilot under tenancy model",
    "Track latency, isolation, and reliability SLAs",
    "Convert to commercial rollout with module licensing",
], color=RGBColor(193, 228, 255), fs=22)

prs.save(str(OUT))
print(f"Created: {OUT}")
