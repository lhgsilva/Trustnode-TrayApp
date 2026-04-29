from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "Trustnode_edge_app" / "docs"
ASSETS = DOCS / "assets"
OUTPUT = DOCS / "TRUSTNODE_SCALABILITY_ARCHITECTURE_DECK_2026-04-22.pptx"
LOGO = ROOT / "trustnode_logo.png"

# Theme colors
NAVY = RGBColor(15, 23, 42)
MID = RGBColor(30, 58, 95)
ACCENT = RGBColor(10, 124, 255)
ACCENT2 = RGBColor(14, 165, 233)
LIGHT_BG = RGBColor(244, 247, 251)
TEXT_DARK = RGBColor(22, 32, 45)
TEXT_MUTED = RGBColor(76, 96, 120)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(21, 128, 61)
ORANGE = RGBColor(180, 83, 9)

prs = Presentation()
# 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.58))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    if title:
        tf = band.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"  {title}"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = WHITE


def add_footer(slide, text="TrustNode - Architecture 2026"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.28), prs.slide_width, Inches(0.28))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(230, 236, 245)
    line.line.fill.background()
    tf = line.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text + "  "
    r.font.size = Pt(10)
    r.font.color.rgb = TEXT_MUTED


def add_logo(slide, x=11.6, y=0.08, h=0.38):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), height=Inches(h))


def add_title_subtitle(slide, title, subtitle):
    t = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.8), Inches(1.2))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(17)
    r2.font.color.rgb = RGBColor(191, 230, 255)


def add_bullets(slide, x, y, w, h, items, font_size=18, dark=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = (TEXT_DARK if dark else WHITE)


def add_section_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.72), Inches(12), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = TEXT_DARK


def add_card(slide, x, y, w, h, title, body, badge=None, badge_color=ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(217, 227, 239)
    card.line.width = Pt(1)

    if badge:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(y + 0.12), Inches(1.2), Inches(0.32))
        b.fill.solid()
        b.fill.fore_color.rgb = badge_color
        b.line.fill.background()
        bt = b.text_frame
        bt.clear()
        bp = bt.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = badge
        br.font.size = Pt(10)
        br.font.bold = True
        br.font.color.rgb = WHITE

    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.52), Inches(w - 0.4), Inches(h - 0.62))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    rr = p.add_run()
    rr.text = title
    rr.font.size = Pt(15)
    rr.font.bold = True
    rr.font.color.rgb = MID

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED


def add_image_full(slide, image_path, x=0.7, y=1.25, w=12.0, h=5.7):
    if image_path.exists():
        # frame
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x-0.06), Inches(y-0.06), Inches(w+0.12), Inches(h+0.12))
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = RGBColor(210, 224, 241)
        frame.line.width = Pt(1)
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_timeline(slide):
    y = 2.0
    phases = [
        ("Phase 0", "1-2 wks", "Pipeline stability\nIngest URL/token diagnostics"),
        ("Phase 1", "2-3 wks", "Control plane core\nTenant/customer/edge/license"),
        ("Phase 2", "2-3 wks", "Tenant routing\nWildcard subdomains"),
        ("Phase 3", "2-4 wks", "Security hardening\nMFA + audits + PITR drills"),
    ]
    x = 0.75
    for i, (ph, dur, txt) in enumerate(phases):
        w = 3.05
        h = 1.52
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + i*3.15), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(204, 220, 240)
        t = shape.text_frame
        t.clear()
        p = t.paragraphs[0]
        r = p.add_run()
        r.text = ph
        r.font.bold = True
        r.font.size = Pt(15)
        r.font.color.rgb = MID
        p2 = t.add_paragraph()
        p2.text = dur
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT
        p3 = t.add_paragraph()
        p3.text = txt
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED


# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, NAVY)
# accent blocks
acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.1), Inches(7.5))
acc.fill.solid(); acc.fill.fore_color.rgb = MID; acc.line.fill.background()
acc2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.65), Inches(13.333), Inches(0.85))
acc2.fill.solid(); acc2.fill.fore_color.rgb = ACCENT; acc2.line.fill.background()
add_logo(slide, x=0.55, y=0.52, h=0.70)
add_title_subtitle(
    slide,
    "TrustNode Scalability\nArchitecture",
    "Low-cost, secure, multi-customer, near-real-time production design",
)
add_bullets(
    slide,
    x=5.2, y=3.2, w=7.4, h=2.4,
    items=[
        "Prepared date: 22 April 2026",
        "Based on current TrustNode edge/cloud architecture",
        "Focus: minimal-change path to commercial rollout",
    ],
    font_size=15,
)

# Slide 2: Executive summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Executive Summary")
add_logo(slide)
add_footer(slide)
add_card(slide, 0.7, 1.2, 4.0, 2.0, "What already works", "Edge-first durability, outbox sync, scoped ingest, tenant-aware schema.", badge="Strong", badge_color=GREEN)
add_card(slide, 4.95, 1.2, 4.0, 2.0, "Main gap", "Customer lifecycle, licensing, and edge activation are not yet first-class.", badge="Gap", badge_color=ORANGE)
add_card(slide, 9.2, 1.2, 3.4, 2.0, "Best approach", "Add a lightweight control plane over existing data plane.", badge="Plan", badge_color=ACCENT)
add_card(slide, 0.7, 3.45, 12.0, 2.9, "Recommendation", "Keep current telemetry architecture. Add control plane + wildcard tenant routing + stricter tenant security controls. This reaches sellable SaaS with minimal implementation risk.", badge="Decision", badge_color=ACCENT2)

# Slide 3: Current architecture diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Current Architecture")
add_logo(slide)
add_footer(slide)
img1 = ASSETS / "topology_1.png"
add_image_full(slide, img1, x=0.65, y=1.0, w=12.1, h=5.9)

# Slide 4: Target architecture diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Target Three-Plane Architecture")
add_logo(slide)
add_footer(slide)
img2 = ASSETS / "topology_2.png"
add_image_full(slide, img2, x=0.65, y=1.0, w=12.1, h=5.9)

# Slide 5: Security topology
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Security & Trust Boundaries")
add_logo(slide)
add_footer(slide)
img3 = ASSETS / "topology_3.png"
add_image_full(slide, img3, x=0.65, y=1.0, w=12.1, h=3.55)
add_card(slide, 0.7, 4.75, 5.95, 2.0, "Critical controls", "Outbound-only edge traffic, separate human/device auth, strict tenant scoping, RLS on user tables.", badge="Must", badge_color=ORANGE)
add_card(slide, 6.85, 4.75, 5.95, 2.0, "Resilience", "PITR backups, ingest audit trail, deterministic retry/idempotency, cloud delay alarms.", badge="Ops", badge_color=ACCENT2)

# Slide 6: Data and tenancy model
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Data Model & Tenant Strategy")
add_logo(slide)
add_footer(slide)
add_card(slide, 0.7, 1.05, 4.0, 2.2, "Raw + Latest split", "telemetry_samples_raw (immutable) + latest_machine_state (live snapshot).", badge="Core", badge_color=ACCENT)
add_card(slide, 4.95, 1.05, 4.0, 2.2, "Tenant isolation", "tenant_id on all rows + API scope checks + DB RLS policies.", badge="Security", badge_color=GREEN)
add_card(slide, 9.2, 1.05, 3.4, 2.2, "Commercial tiers", "Tier A shared DB, Tier B isolated schema, Tier C dedicated stack.", badge="Scale", badge_color=ACCENT2)

# table-like layout using shapes
hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.55), Inches(12.0), Inches(0.45))
hdr.fill.solid(); hdr.fill.fore_color.rgb = MID; hdr.line.fill.background()
ht = hdr.text_frame; ht.clear(); hp = ht.paragraphs[0]; hr = hp.add_run(); hr.text = "Recommended rollout by customer tier"; hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = WHITE
rows = [
    ("Tier A (default)", "Shared DB + RLS", "Low", "SMB / standard deployments"),
    ("Tier B", "Shared cluster + dedicated schema", "Medium", "Regulated mid-size customers"),
    ("Tier C", "Dedicated project/cluster", "High", "Enterprise/compliance-heavy"),
]
for i, row in enumerate(rows):
    y = 4.03 + i*0.68
    for j, text in enumerate(row):
        x = [0.7, 3.55, 7.9, 9.35][j]
        w = [2.85, 4.35, 1.35, 3.35][j]
        c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.66))
        c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.color.rgb = RGBColor(214, 225, 240)
        tf = c.text_frame; tf.clear(); p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(11); p.font.color.rgb = TEXT_DARK

# Slide 7: URL and customer experience
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Customer URL, Access & Licensing")
add_logo(slide)
add_footer(slide)
add_card(slide, 0.7, 1.1, 6.0, 2.5, "URL strategy", "Use wildcard subdomains: customer-a.trustnode.lsapps.app. Resolve host header -> tenant context in backend.", badge="Routing", badge_color=ACCENT)
add_card(slide, 7.0, 1.1, 5.7, 2.5, "Licensing", "Control modules by tenant license and role claims. Same user can access only allowed modules.", badge="License", badge_color=GREEN)
add_card(slide, 0.7, 3.9, 12.0, 2.8, "Control plane entities", "tenants, customers, edges, licenses, license_modules, edge_activation_codes, user_tenant_memberships, password_reset_events, security_audit_log", badge="Schema", badge_color=ACCENT2)

# Slide 8: Rollout plan
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Implementation Timeline")
add_logo(slide)
add_footer(slide)
add_timeline(slide)

# Slide 9: SLAs and monitoring
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "SLA Targets & Observability")
add_logo(slide)
add_footer(slide)
add_card(slide, 0.7, 1.1, 3.8, 2.0, "Local live latency", "<= 1 second from sample commit to local UI", badge="SLA", badge_color=ACCENT)
add_card(slide, 4.8, 1.1, 3.8, 2.0, "Cloud live latency", "<= 2 seconds p95 from edge sample_ts to cloud UI", badge="SLA", badge_color=ACCENT)
add_card(slide, 8.9, 1.1, 3.8, 2.0, "Tenant isolation", "0 cross-tenant data leakage incidents", badge="Security", badge_color=GREEN)
add_card(slide, 0.7, 3.4, 12.0, 3.1, "Key metrics to monitor", "PLC read duration, collection loop duration, local DB write latency, outbox depth, oldest unsynced age, batch upload latency, ack success rate, duplicate rate, rejected rate, websocket clients, end-to-end cloud latency.", badge="Metrics", badge_color=ACCENT2)

# Slide 10: Risks and mitigations
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, LIGHT_BG)
add_top_band(slide, "Top Risks & Mitigations")
add_logo(slide)
add_footer(slide)
add_card(slide, 0.7, 1.1, 6.0, 1.7, "Risk 1: Sync misconfiguration", "Mitigation: startup diagnostics + hard fail warnings for missing ingest URL/token.", badge="High", badge_color=ORANGE)
add_card(slide, 7.0, 1.1, 5.7, 1.7, "Risk 2: Tenant bleed", "Mitigation: host->tenant binding, token scope checks, RLS policy tests in CI.", badge="Critical", badge_color=ORANGE)
add_card(slide, 0.7, 3.0, 6.0, 1.7, "Risk 3: Backlog growth", "Mitigation: adaptive batching, retry jitter, backlog alarms, capacity runbooks.", badge="High", badge_color=ORANGE)
add_card(slide, 7.0, 3.0, 5.7, 1.7, "Risk 4: Admin account compromise", "Mitigation: enforce MFA + privileged audit logs + password reset policy.", badge="Critical", badge_color=ORANGE)
add_card(slide, 0.7, 4.9, 12.0, 1.8, "Go-live gate", "No production rollout without cross-tenant isolation tests, restore drill pass, and p95 cloud latency validation.", badge="Gate", badge_color=ACCENT2)

# Slide 11: Next steps
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, NAVY)
add_logo(slide, x=11.7, y=0.08, h=0.38)
add_top_band(slide, "Execution Next Steps")
add_footer(slide, text="TrustNode | Execute Phase 0 first")

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(9.2), Inches(0.9))
tf = title_box.text_frame
tf.clear()
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Recommended 30 / 60 / 90 Day Plan"
r.font.size = Pt(31)
r.font.bold = True
r.font.color.rgb = WHITE

add_bullets(
    slide,
    x=0.9, y=2.25, w=12.0, h=4.6,
    items=[
        "30 days: stabilize ingest bootstrap, add diagnostics, add sync alarms.",
        "60 days: deploy control plane core (tenant/customer/edge/license + activation).",
        "90 days: enforce MFA, wildcard customer routing, full isolation test suite.",
        "Outcome: sellable, secure, low-cost multi-customer platform with clean scale path.",
    ],
    font_size=19,
)

prs.save(str(OUTPUT))
print(f"Created: {OUTPUT}")
