"""Build a 22-slide pitch deck for the TrustNode security & architecture whitepaper.

Output:
  docs/TrustNode_Security_Deck.pptx

Slide plan:
   1. Cover
   2. The problem we solve
   3. What TrustNode is, in one sentence
   4. Architecture — single-customer deployment       [diagram]
   5. Purdue model alignment                          [diagram]
   6. Read-only on the plant floor (protocols)
   7. Outbound-only network — no inbound to the edge
   8. Three login roles, enforced server-side         [diagram]
   9. Multi-tenant isolation                          [diagram]
  10. Authentication & data protection
  11. Store-and-forward resilience                    [diagram]
  12. Deployment options — overview
  13. Deployment A — Plant PC                         [diagram]
  14. Deployment B — Industrial PC in panel           [diagram]
  15. Deployment C — Customer server                  [diagram]
  16. Deployment D — Cloud-bridged                    [diagram]
  17. Deployment E — Multi-plant central              [diagram]
  18. Storage options — five ways                     [diagram]
  19. Networking prerequisites
  20. Backup & redundancy
  21. Compliance map (IEC 62443 / NIS2 / ISO 27001 / Purdue)
  22. Why customers can trust us / Q&A
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

# -------------------------- brand colors (RGBColor uses 0..255 triples) -----
NAVY = RGBColor(0x0E, 0x1A, 0x3A)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
SLATE = RGBColor(0x2B, 0x35, 0x48)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF4, 0xF6, 0xFA)
INK = RGBColor(0x0E, 0x11, 0x16)
MUTED = RGBColor(0x5B, 0x64, 0x73)
ALERT = RGBColor(0xC2, 0x5B, 0x35)
GOOD = RGBColor(0x2D, 0x7A, 0x4F)
GRID = RGBColor(0xD7, 0xDC, 0xE5)
CALLOUT_BG = RGBColor(0xEE, 0xF3, 0xFB)

# ------------------------------------------------------------ paths --------
ROOT = Path(__file__).resolve().parents[1]
DIAG_DIR = ROOT / "docs" / "diagrams"
OUT_PATH = ROOT / "docs" / "TrustNode_Security_Deck.pptx"

# Slide size (16:9 widescreen)
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

TOTAL_SLIDES = 22

# ---------------------------------------- helpers ---------------------------
def fill_rect(slide, left, top, w, h, color, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, w, h, text, *, color=INK, size=18, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0); tf.margin_right = Cm(0)
    tf.margin_top = Cm(0); tf.margin_bottom = Cm(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, w, h, items, *, color=INK, size=16,
                bullet_color=TEAL, anchor=MSO_ANCHOR.TOP, after=6):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Cm(0); tf.margin_top = Cm(0); tf.margin_bottom = Cm(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        rb = p.add_run()
        rb.text = "•  "
        rb.font.name = "Calibri"; rb.font.size = Pt(size); rb.font.bold = True
        rb.font.color.rgb = bullet_color
        rt = p.add_run()
        rt.text = item
        rt.font.name = "Calibri"; rt.font.size = Pt(size)
        rt.font.color.rgb = color
        p.space_after = Pt(after)


def add_header(slide, title, subtitle=None):
    fill_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(2.6), NAVY)
    fill_rect(slide, Cm(0), Cm(2.6), SLIDE_W, Cm(0.18), TEAL)
    add_text(slide, Cm(1.6), Cm(0.4), Cm(SLIDE_W.cm - 3.0), Cm(2.2),
             title, color=PAPER, size=30, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Cm(1.6), Cm(2.95), Cm(SLIDE_W.cm - 3.0), Cm(1.0),
                 subtitle, color=MUTED, size=14)


def add_footer(slide, page_num):
    fill_rect(slide, Cm(0), Cm(SLIDE_H.cm - 0.8), SLIDE_W, Cm(0.04), GRID)
    add_text(slide, Cm(1.6), Cm(SLIDE_H.cm - 0.7), Cm(20), Cm(0.6),
             "TrustNode Edge — Security & Architecture", color=MUTED, size=10)
    add_text(slide, Cm(SLIDE_W.cm - 4), Cm(SLIDE_H.cm - 0.7), Cm(2.4), Cm(0.6),
             f"{page_num} / {TOTAL_SLIDES}", color=MUTED, size=10, align=PP_ALIGN.RIGHT)


def add_diagram(slide, png_name, *, left=Cm(1.6), top=Cm(4.2),
                w=Cm(SLIDE_W.cm - 3.2), h=Cm(13.2)):
    p = DIAG_DIR / png_name
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width=w, height=h)


def diagram_with_caption(slide, png, caption_left, caption_text):
    """Diagram on left 2/3, caption on right 1/3 — good for topology slides."""
    add_diagram(slide, png, left=Cm(1.0), top=Cm(4.2), w=Cm(20.0), h=Cm(13.2))
    fill_rect(slide, Cm(22.2), Cm(4.2), Cm(10.4), Cm(13.2), SOFT, line=GRID)
    fill_rect(slide, Cm(22.2), Cm(4.2), Cm(0.2), Cm(13.2), TEAL)
    add_text(slide, Cm(22.6), Cm(4.5), Cm(9.8), Cm(1.4),
             caption_left, color=NAVY, size=20, bold=True)
    add_text(slide, Cm(22.6), Cm(6.2), Cm(9.8), Cm(11.0),
             caption_text, color=INK, size=13)


# ============================================================ deck builder
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # ----------------------------------- slide 1 — Cover --------------------
    s = prs.slides.add_slide(blank)
    fill_rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, NAVY)
    fill_rect(s, Cm(0), Cm(11.4), SLIDE_W, Cm(0.25), TEAL)
    add_text(s, Cm(2.0), Cm(4.4), Cm(SLIDE_W.cm - 4), Cm(3.2),
             "TrustNode Edge", color=PAPER, size=64, bold=True)
    add_text(s, Cm(2.0), Cm(7.4), Cm(SLIDE_W.cm - 4), Cm(2.0),
             "Security & Architecture", color=TEAL, size=36, bold=True)
    add_text(s, Cm(2.0), Cm(9.8), Cm(SLIDE_W.cm - 4), Cm(1.4),
             "Built for industrial operators. Read-only on the plant floor.\n"
             "Outbound-only on the network. Tenant-isolated end-to-end.",
             color=RGBColor(0xCD, 0xD5, 0xE0), size=18)
    add_text(s, Cm(2.0), Cm(15.5), Cm(SLIDE_W.cm - 4), Cm(2.5),
             "Audience: Plant Managers · IT Security · OT Engineering · Compliance\n"
             "Document version: 2026-05-15 (rev. 2)",
             color=RGBColor(0x9A, 0xA3, 0xB3), size=12)

    # ----------------------------------- slide 2 — Problem ------------------
    s = prs.slides.add_slide(blank)
    add_header(s, "The problem we solve",
               "Industrial customers want their data — without changing the plant.")
    add_bullets(s, Cm(2.0), Cm(4.6), Cm(SLIDE_W.cm - 4), Cm(8), [
        "Plants generate enormous amounts of data, locked inside PLCs and meters.",
        "Most monitoring tools require writing back to the PLC or opening inbound ports — neither is acceptable.",
        "IT teams must protect the corporate network from anything talking to the plant floor.",
        "OT teams must protect production from anything that might disturb control logic.",
        "Compliance teams need an audit trail, tenant isolation, and a story for IEC 62443 / NIS2.",
    ], size=20)
    add_footer(s, 2)

    # ----------------------------------- slide 3 — What we are -------------
    s = prs.slides.add_slide(blank)
    add_header(s, "What TrustNode is, in one sentence",
               "A read-only industrial data gateway with a cloud portal that never reaches back.")
    fill_rect(s, Cm(2.6), Cm(5.4), Cm(SLIDE_W.cm - 5.2), Cm(7.5), SOFT, line=GRID)
    fill_rect(s, Cm(2.6), Cm(5.4), Cm(0.3), Cm(7.5), TEAL)
    add_text(s, Cm(3.6), Cm(6.2), Cm(SLIDE_W.cm - 7.4), Cm(5.5),
             "PLC  →  Edge  →  Cloud  →  Browser",
             color=NAVY, size=44, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Cm(3.6), Cm(10.6), Cm(SLIDE_W.cm - 7.4), Cm(2.0),
             "Always that direction. The cloud cannot reach the edge.\n"
             "The browser cannot reach the PLC.",
             color=MUTED, size=18, align=PP_ALIGN.CENTER)
    add_footer(s, 3)

    # ----------------------------------- slide 4 — Single customer ---------
    s = prs.slides.add_slide(blank)
    add_header(s, "Single-customer deployment",
               "Plant floor → edge → cloud → browsers. Only outbound HTTPS crosses the boundary.")
    add_diagram(s, "architecture_single_customer.png")
    add_footer(s, 4)

    # ----------------------------------- slide 5 — Purdue ------------------
    s = prs.slides.add_slide(blank)
    add_header(s, "Purdue model alignment",
               "Where TrustNode sits inside the industry's reference picture of plant networks.")
    add_diagram(s, "architecture_purdue.png")
    add_footer(s, 5)

    # ----------------------------------- slide 6 — Protocols --------------
    s = prs.slides.add_slide(blank)
    add_header(s, "Read-only on the plant floor",
               "Four industrial protocols. Zero write code paths.")
    cards = [
        ("OPC UA  (IEC 62541)",
         "Vendor-neutral standard. Siemens, Beckhoff, generic OPC servers.\nRead-only tag access; server's auth & encryption respected."),
        ("Siemens S7 (Snap7)",
         "Native S7-300/400/1200/1500. Read-only DB/M/I/Q access.\nPLC's own ACL still applies."),
        ("Allen-Bradley (EtherNet/IP)",
         "CompactLogix / ControlLogix via pycomm3 / pylogix.\nRead-only. No write_tag bindings anywhere."),
        ("Modbus TCP",
         "Universal protocol for power meters, VFDs, sensors.\nReads only holding & input registers."),
    ]
    for idx, (title, body) in enumerate(cards):
        col = idx % 2; row = idx // 2
        left = Cm(1.6 + col * 15.6); top = Cm(4.6 + row * 6.4)
        fill_rect(s, left, top, Cm(15.0), Cm(5.8), PAPER, line=GRID)
        fill_rect(s, left, top, Cm(15.0), Cm(0.18), TEAL)
        add_text(s, left + Cm(0.6), top + Cm(0.4), Cm(14), Cm(1.2),
                 title, color=NAVY, size=18, bold=True)
        add_text(s, left + Cm(0.6), top + Cm(1.8), Cm(14), Cm(4),
                 body, color=INK, size=14)
    add_footer(s, 6)

    # ----------------------------------- slide 7 — Outbound-only ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Outbound-only network",
               "No inbound port on the edge, no inbound from cloud to plant.")
    add_bullets(s, Cm(1.6), Cm(4.6), Cm(15.4), Cm(13.0), [
        "Edge dials out to cloud over HTTPS / TCP 443.",
        "Same direction as a web browser fetching a page.",
        "No NAT rule, no port forwarding, no listener exposed.",
        "Cloud VPS exposes only TCP/443 through nginx.",
        "FastAPI backend binds 127.0.0.1 — only reachable via nginx.",
        "No SSH server, no RDP listener, no reverse tunnel on edge.",
    ], size=18)
    fill_rect(s, Cm(18.0), Cm(4.6), Cm(13.8), Cm(13), SOFT, line=GRID)
    fill_rect(s, Cm(18.0), Cm(4.6), Cm(0.3), Cm(13), TEAL)
    add_text(s, Cm(18.6), Cm(5.0), Cm(12.8), Cm(1.0),
             "Why this matters", color=TEAL, size=18, bold=True)
    add_text(s, Cm(18.6), Cm(6.4), Cm(12.8), Cm(11.0),
             "Two of the largest industrial-breach categories of the past decade were "
             "remote-desktop tools left exposed and SSH servers without strong keys.\n\n"
             "TrustNode simply does not run either.\n\n"
             "Your perimeter stays closed. If your firewall is healthy today, "
             "installing TrustNode does not change that.",
             color=INK, size=15)
    add_footer(s, 7)

    # ----------------------------------- slide 8 — Three roles -----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Three login roles, enforced server-side",
               "Same login endpoint, three very different views.")
    add_diagram(s, "architecture_three_role.png")
    add_footer(s, 8)

    # ----------------------------------- slide 9 — Multi-tenant ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Multi-tenant isolation",
               "One VPS, many customers, zero cross-tenant data.")
    add_diagram(s, "architecture_multi_tenant.png")
    add_footer(s, 9)

    # ----------------------------------- slide 10 — Auth + data ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Authentication & data protection",
               "Industry-standard cryptography, plain to explain to your CISO.")
    fill_rect(s, Cm(1.6), Cm(4.6), Cm(15.0), Cm(13.0), PAPER, line=GRID)
    fill_rect(s, Cm(1.6), Cm(4.6), Cm(15.0), Cm(0.18), TEAL)
    add_text(s, Cm(2.2), Cm(4.9), Cm(14), Cm(1.0), "Authentication", color=NAVY, size=20, bold=True)
    add_bullets(s, Cm(2.2), Cm(6.1), Cm(14), Cm(11), [
        "JWT (RFC 7519) with HS256 signatures.",
        "12-hour token lifetime; no refresh tokens.",
        "Passwords stored as PBKDF2-HMAC-SHA256, 120,000 iterations.",
        "Per-user salt; no plaintext ever logged or persisted.",
        "Three roles: master / customer-admin / client viewer.",
        "Every API call re-validates the JWT.",
    ], size=14)
    fill_rect(s, Cm(17.2), Cm(4.6), Cm(15.0), Cm(13.0), PAPER, line=GRID)
    fill_rect(s, Cm(17.2), Cm(4.6), Cm(15.0), Cm(0.18), TEAL)
    add_text(s, Cm(17.8), Cm(4.9), Cm(14), Cm(1.0), "Data protection", color=NAVY, size=20, bold=True)
    add_bullets(s, Cm(17.8), Cm(6.1), Cm(14), Cm(11), [
        "TLS 1.2+ everywhere — browsers, edge↔cloud, DB↔backend.",
        "Postgres / Supabase: sslmode=require, provider encryption at rest.",
        "WSS auto-upgrade — mixed-content is structurally impossible.",
        "Edge SQLite on BitLocker / LUKS / FileVault for disk encryption.",
        "Activation codes: SHA-256 hashed, 30-min TTL, one-shot.",
        "Password reset tokens: hashed, 15-min TTL, one-shot.",
    ], size=14)
    add_footer(s, 10)

    # ----------------------------------- slide 11 — Store-and-forward ----
    s = prs.slides.add_slide(blank)
    add_header(s, "Store-and-forward resilience",
               "Cloud outages do not become plant-floor data gaps.")
    add_diagram(s, "architecture_store_forward.png")
    add_footer(s, 11)

    # ============================ NEW SLIDES ============================
    # ----------------------------------- slide 12 — Deployment overview --
    s = prs.slides.add_slide(blank)
    add_header(s, "Deployment options — overview",
               "Same software, five common shapes. Pick what fits your IT/OT reality.")
    summary = [
        ("A — Plant PC",       "Windows PC inside the plant. Local-only.",       "Smallest sites, pilots, no cloud needed."),
        ("B — IPC in panel",   "Rugged DIN-rail PC in the electrical cabinet.",  "Machine builders, OEMs, no server room."),
        ("C — Customer server","Linux VM in the customer's datacenter.",         "Mid/large customers with their own IT."),
        ("D — Cloud-bridged",  "Edge at plant + cloud portal. The default.",     "Web access from anywhere; multi-site."),
        ("E — Multi-plant",    "Several plants pooled into one central historian.", "Multi-site enterprises, one source of truth."),
    ]
    for i, (title, body, best) in enumerate(summary):
        top = Cm(4.6 + i * 2.7)
        fill_rect(s, Cm(1.6), top, Cm(SLIDE_W.cm - 3.2), Cm(2.4), PAPER, line=GRID)
        fill_rect(s, Cm(1.6), top, Cm(0.18), Cm(2.4), TEAL)
        add_text(s, Cm(2.0), top + Cm(0.2), Cm(7), Cm(2),
                 title, color=NAVY, size=16, bold=True)
        add_text(s, Cm(9.5), top + Cm(0.2), Cm(13), Cm(2),
                 body, color=INK, size=14)
        add_text(s, Cm(23), top + Cm(0.2), Cm(SLIDE_W.cm - 24.5), Cm(2),
                 "Best for: " + best, color=MUTED, size=12)
    add_footer(s, 12)

    # ----------------------------------- slide 13 — Topology A ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Deployment A — Plant PC / desktop install",
               "Simplest setup. TrustNode runs on a Windows PC inside the plant.")
    diagram_with_caption(
        s, "deployment_plant_pc.png",
        "Plant PC",
        "• Windows 10/11 PC inside the plant LAN.\n"
        "• Local SQLite store, no cloud needed.\n"
        "• Operator browser hits it on the LAN.\n"
        "• 16 GB RAM / 256 GB SSD typical.\n\n"
        "Air-gap friendly. Best for small standalone sites and pilots.",
    )
    add_footer(s, 13)

    # ----------------------------------- slide 14 — Topology B ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Deployment B — Industrial PC (IPC) in the panel",
               "Rugged DIN-rail PC, fanless, 24 VDC. Mounted next to the PLCs.")
    diagram_with_caption(
        s, "deployment_ipc_panel.png",
        "IPC in panel",
        "• Fanless, DIN-rail, 24 VDC supply.\n"
        "• 8–16 GB RAM, industrial SSD/SD.\n"
        "• Common: Siemens IPC127E, B&R APC910, Beckhoff CX, Advantech UNO.\n"
        "• Optional cellular / VPN modem in same panel for cloud portal.\n\n"
        "Best for machine builders shipping a turnkey line.",
    )
    add_footer(s, 14)

    # ----------------------------------- slide 15 — Topology C ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Deployment C — Customer server in datacenter",
               "Linux VM or rack server. Often serves several plants from one box.")
    diagram_with_caption(
        s, "deployment_customer_server.png",
        "Customer server",
        "• 4 vCPU / 8 GB RAM (start), Ubuntu / RHEL / Debian.\n"
        "• Reaches each plant's PLCs over the customer's WAN/VPN.\n"
        "• Writes into the customer's own PostgreSQL.\n"
        "• Customer's IT runs backups, monitoring, HA.\n\n"
        "Best for mid/large customers with internal IT and on-prem policy.",
    )
    add_footer(s, 15)

    # ----------------------------------- slide 16 — Topology D ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Deployment D — Cloud-bridged (the default)",
               "Edge in the plant + cloud portal. Web access from anywhere.")
    diagram_with_caption(
        s, "deployment_cloud_bridged.png",
        "Cloud-bridged",
        "• Edge can be PC, IPC, or VM (A/B/C).\n"
        "• Outbound HTTPS to a cloud VPS — ours or yours.\n"
        "• Cloud handles dashboards, multi-tenant portal, audit log.\n"
        "• Edge keeps SQLite buffer (store-and-forward).\n\n"
        "Best when web access from anywhere matters.",
    )
    add_footer(s, 16)

    # ----------------------------------- slide 17 — Topology E ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Deployment E — Multi-plant central historian",
               "Several plants, one customer, one central historian.")
    diagram_with_caption(
        s, "deployment_multi_plant.png",
        "Multi-plant",
        "• One edge per plant; one central server (on-prem or cloud).\n"
        "• tenant_id = customer; edge_id distinguishes plants.\n"
        "• All plants visible in one portal, filterable per plant.\n"
        "• PostgreSQL on customer infra OR managed cloud DB.\n\n"
        "Best for multi-site enterprises wanting one source of truth.",
    )
    add_footer(s, 17)

    # ----------------------------------- slide 18 — Storage options -----
    s = prs.slides.add_slide(blank)
    add_header(s, "Storage options — where the data physically lives",
               "Five options, pick what fits IT policy and budget.")
    add_diagram(s, "storage_options.png")
    add_footer(s, 18)

    # ----------------------------------- slide 19 — Networking ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Networking prerequisites",
               "The checklist your IT and OT teams can run through together.")
    # Two columns: must reach / does NOT need
    fill_rect(s, Cm(1.6), Cm(4.6), Cm(15.0), Cm(13.0), PAPER, line=GRID)
    fill_rect(s, Cm(1.6), Cm(4.6), Cm(15.0), Cm(0.18), GOOD)
    add_text(s, Cm(2.2), Cm(4.9), Cm(14), Cm(1.0),
             "TrustNode MUST reach", color=NAVY, size=20, bold=True)
    add_bullets(s, Cm(2.2), Cm(6.1), Cm(14), Cm(11), [
        "PLCs at configured IPs (typical ports: 502 Modbus, 4840 OPC UA, 102 S7, 44818 EtherNet/IP).",
        "DNS resolution (or use IPs).",
        "NTP for time sync — historian only useful if timestamps are correct.",
        "Outbound HTTPS/443 to cloud (cloud-bridged topologies only).",
        "Customer's proxy / custom CA bundle if enforced.",
    ], size=13)

    fill_rect(s, Cm(17.2), Cm(4.6), Cm(15.0), Cm(13.0), PAPER, line=GRID)
    fill_rect(s, Cm(17.2), Cm(4.6), Cm(15.0), Cm(0.18), ALERT)
    add_text(s, Cm(17.8), Cm(4.9), Cm(14), Cm(1.0),
             "TrustNode does NOT need", color=NAVY, size=20, bold=True)
    add_bullets(s, Cm(17.8), Cm(6.1), Cm(14), Cm(11), [
        "Inbound ports on the edge — none.",
        "Internet access from the PLCs themselves.",
        "Open inbound firewall to the cloud.",
        "RDP / SSH / VNC servers on the edge.",
        "Domain admin or service-account privileges.",
    ], size=13, bullet_color=ALERT)
    add_footer(s, 19)

    # ----------------------------------- slide 20 — Backup --------------
    s = prs.slides.add_slide(blank)
    add_header(s, "Backup & redundancy",
               "Who runs the backups depends on which storage option you picked.")
    rows = [
        ("Local SQLite only",       "Customer",       "Periodic file copy to NAS / external disk."),
        ("Customer Postgres",       "Customer IT",    "Postgres replication + customer's standard DB ops."),
        ("Managed cloud DB",        "Provider",       "Daily backups, point-in-time restore, multi-AZ replicas."),
        ("Hybrid (edge + cloud)",   "Both",           "Cloud backups + local SQLite is permanent backup-of-record."),
        ("Other DB on request",     "Varies",         "Whatever that DB engine offers natively."),
    ]
    # header
    fill_rect(s, Cm(1.6), Cm(4.6), Cm(SLIDE_W.cm - 3.2), Cm(1.3), NAVY)
    add_text(s, Cm(2.0), Cm(4.85), Cm(9), Cm(0.9),
             "Storage option", color=PAPER, size=14, bold=True)
    add_text(s, Cm(12.0), Cm(4.85), Cm(7), Cm(0.9),
             "Backup owner", color=PAPER, size=14, bold=True)
    add_text(s, Cm(19.5), Cm(4.85), Cm(SLIDE_W.cm - 21), Cm(0.9),
             "Redundancy options", color=PAPER, size=14, bold=True)
    for i, (col1, col2, col3) in enumerate(rows):
        top = Cm(6.0 + i * 1.5)
        alt = (i % 2 == 1)
        fill_rect(s, Cm(1.6), top, Cm(SLIDE_W.cm - 3.2), Cm(1.4),
                 RGBColor(0xF9, 0xFA, 0xFC) if alt else PAPER, line=GRID)
        add_text(s, Cm(2.0), top + Cm(0.15), Cm(9), Cm(1.1),
                 col1, color=NAVY, size=13, bold=True)
        add_text(s, Cm(12.0), top + Cm(0.15), Cm(7), Cm(1.1),
                 col2, color=INK, size=13)
        add_text(s, Cm(19.5), top + Cm(0.15), Cm(SLIDE_W.cm - 21), Cm(1.1),
                 col3, color=INK, size=12)
    add_text(s, Cm(1.6), Cm(14.4), Cm(SLIDE_W.cm - 3.2), Cm(1.2),
             "Edge SQLite always retains a recent buffer regardless of long-term storage choice.",
             color=GOOD, size=14, bold=True)
    add_text(s, Cm(1.6), Cm(15.6), Cm(SLIDE_W.cm - 3.2), Cm(1.6),
             "HA options on request: active-passive edge pair · load-balanced FastAPI · Postgres HA mode.",
             color=MUTED, size=12)
    add_footer(s, 20)

    # ----------------------------------- slide 21 — Compliance ----------
    s = prs.slides.add_slide(blank)
    add_header(s, "Compliance map",
               "We are not certified; we provide the evidence package.")
    standards = [
        ("IEC 62443",
         "International standard for industrial cybersecurity.\n"
         "Zone & conduit (3-2), system security requirements (3-3)."),
        ("NIS2",
         "EU directive on critical-infrastructure cybersecurity.\n"
         "Risk management, incident handling, business continuity (Art. 21)."),
        ("ISO 27001",
         "General information-security management standard.\n"
         "Annex A — access, cryptography, supply chain, logging."),
        ("Purdue Model",
         "OT reference architecture for plant networks.\n"
         "TrustNode at L3 (Operations); cloud at L4–L5."),
    ]
    for idx, (title, body) in enumerate(standards):
        col = idx % 2; row = idx // 2
        left = Cm(1.6 + col * 15.6); top = Cm(4.6 + row * 6.4)
        fill_rect(s, left, top, Cm(15.0), Cm(5.8), PAPER, line=GRID)
        fill_rect(s, left, top, Cm(15.0), Cm(0.18), TEAL)
        add_text(s, left + Cm(0.6), top + Cm(0.4), Cm(14), Cm(1.2),
                 title, color=NAVY, size=20, bold=True)
        add_text(s, left + Cm(0.6), top + Cm(1.8), Cm(14), Cm(4),
                 body, color=INK, size=14)
    add_footer(s, 21)

    # ----------------------------------- slide 22 — Closing -------------
    s = prs.slides.add_slide(blank)
    add_header(s, "Why customers can trust us",
               "Five reasons to keep on the prospect's notebook.")
    add_bullets(s, Cm(1.6), Cm(4.6), Cm(SLIDE_W.cm - 3.2), Cm(11), [
        "Read-only on the plant floor. No PLC writes anywhere in the gateway code.",
        "Outbound-only on the network. Same direction as web traffic, no inbound to the edge.",
        "One customer cannot see another. 49+ automated checks per release verify this.",
        "Your data, your database. Five storage options. Migrate at any time.",
        "Continuous hardening. Audit log is append-only, dependencies are pinned, smokes run on every release.",
    ], size=20)
    fill_rect(s, Cm(1.6), Cm(15.8), Cm(SLIDE_W.cm - 3.2), Cm(1.2), TEAL)
    add_text(s, Cm(2.0), Cm(15.85), Cm(SLIDE_W.cm - 4), Cm(1.0),
             "Questions? Walk us through the full whitepaper at "
             "docs/TRUSTNODE_SECURITY_AND_ARCHITECTURE_WHITEPAPER.md",
             color=PAPER, size=14, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 22)

    prs.save(str(OUT_PATH))
    print(OUT_PATH)


if __name__ == "__main__":
    build()
